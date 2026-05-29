#!/usr/bin/env python3
"""Extract text from PPTX slides into a stable, reviewable text file."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path


def iter_shape_text(shape):
    if getattr(shape, "has_text_frame", False) and shape.text:
        yield shape.text
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                yield " | ".join(cells)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from iter_shape_text(child)


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX text slide by slide.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx is required: python3 -m pip install python-pptx", file=sys.stderr)
        return 2

    prs = Presentation(str(args.pptx))
    args.output.parent.mkdir(parents=True, exist_ok=True)

    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {idx} ===")
        seen: set[str] = set()
        for shape in slide.shapes:
            for text in iter_shape_text(shape):
                cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
                if cleaned and cleaned not in seen:
                    lines.append(cleaned)
                    seen.add(cleaned)
        lines.append("")

    args.output.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
